#!/usr/bin/env python3
"""Faithful-enough pptx -> png rasterizer for QA: reads shapes/positions/fills/text/pictures from the
ACTUAL saved deck and draws them. Not a full renderer, but reflects layout/overlap/overflow truthfully."""
import sys, io
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPI = 120; EMU = 914400
def px(v): return int(v * PPI / EMU)
FONT = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONTB = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
def font(sz, bold): return ImageFont.truetype(FONTB if bold else FONT, max(8, int(sz * PPI / 72)))

def rgb(c):
    try: return (c[0], c[1], c[2])
    except Exception: return None
def fill_rgb(sp):
    try:
        if sp.fill.type is not None: return rgb(sp.fill.fore_color.rgb)
    except Exception: pass
    return None
def line_rgb(sp):
    try: return rgb(sp.line.color.rgb)
    except Exception: return None

def wrap(draw, text, fnt, maxw):
    out = []
    for para in text.split("\n"):
        words = para.split(" "); line = ""
        for w in words:
            t = (line + " " + w).strip()
            if draw.textlength(t, font=fnt) <= maxw or not line: line = t
            else: out.append(line); line = w
        out.append(line)
    return out

def draw_text(draw, box, tf, center=False):
    x0, y0, x1, y1 = box
    # collect paragraphs -> (text, size, color, bold)
    paras = []
    for p in tf.paragraphs:
        t = "".join(r.text for r in p.runs)
        if not t: continue
        r0 = p.runs[0]; sz = r0.font.size.pt if r0.font.size else 14
        col = rgb(r0.font.color.rgb) if (r0.font.color and r0.font.color.type is not None) else (230, 230, 230)
        paras.append((t, sz, col or (230, 230, 230), bool(r0.font.bold)))
    # total height
    lines = []
    for t, sz, col, b in paras:
        fnt = font(sz, b)
        for ln in wrap(draw, t, fnt, x1 - x0):
            lines.append((ln, fnt, col, sz))
    th = sum(int(sz * PPI / 72 * 1.25) for _, _, _, sz in lines)
    y = y0 + (max(0, (y1 - y0) - th) // 2 if center else 4)
    for ln, fnt, col, sz in lines:
        w = draw.textlength(ln, font=fnt)
        x = x0 + (( (x1 - x0) - w) // 2 if center else 4)
        draw.text((x, y), ln, font=fnt, fill=col)
        y += int(sz * PPI / 72 * 1.25)

def render(prs, idx, slide):
    W = px(prs.slide_width); H = px(prs.slide_height)
    img = Image.new("RGB", (W, H), (30, 30, 34)); d = ImageDraw.Draw(img)
    for sp in slide.shapes:
        try: x, y, w, h = px(sp.left), px(sp.top), px(sp.width), px(sp.height)
        except Exception: continue
        if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                im = Image.open(io.BytesIO(sp.image.blob)).convert("RGB").resize((max(1, w), max(1, h)))
                img.paste(im, (x, y))
            except Exception: pass
            continue
        f = fill_rgb(sp); l = line_rgb(sp)
        if f or l:
            box = [x, y, x + w, y + h]
            rad = min(18, w // 4, h // 4)
            try: d.rounded_rectangle(box, radius=max(0, rad), fill=f, outline=l, width=2 if l else 0)
            except Exception: d.rectangle(box, fill=f, outline=l)
        if sp.has_text_frame and sp.text_frame.text.strip():
            draw_text(d, (x, y, x + w, y + h), sp.text_frame, center=(f is not None))
    return img

prs = Presentation(sys.argv[1] if len(sys.argv) > 1 else "Dog_Fur_Three_Approaches.pptx")
imgs = [render(prs, i, s) for i, s in enumerate(prs.slides)]
for i, im in enumerate(imgs): im.save(f"/tmp/slide_{i+1}.png")
# contact sheet 2 cols
import math
cols = 2; rows = math.ceil(len(imgs) / cols); tw = imgs[0].width // 2; th = imgs[0].height // 2
sheet = Image.new("RGB", (tw * cols + 20, th * rows + 20 * rows), (10, 10, 12))
for i, im in enumerate(imgs):
    t = im.resize((tw, th)); r, c = divmod(i, cols)
    sheet.paste(t, (c * tw + 10, r * (th + 15) + 10))
sheet.save("/tmp/deck_contact.png")
print(f"rendered {len(imgs)} slides -> /tmp/slide_*.png + /tmp/deck_contact.png")
