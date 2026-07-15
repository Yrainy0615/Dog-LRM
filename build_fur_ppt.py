#!/usr/bin/env python3
"""Build the three-approach fur comparison deck (16:9). Flow diagrams = native shapes; result images
embedded; comparison table native. Run in dog-lrm env (python-pptx)."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

W, H = Inches(13.333), Inches(7.5)
BG   = RGBColor(0x1E, 0x1E, 0x22)
FG   = RGBColor(0xEC, 0xEC, 0xEC)
SUB  = RGBColor(0xA8, 0xA8, 0xB0)
BLUE = RGBColor(0x3B, 0x82, 0xF6)   # A feed-forward
RED  = RGBColor(0x8A, 0x50, 0x50)   # B two-stage (dead end)
GREEN= RGBColor(0x22, 0xA5, 0x6A)   # C cascade
PURP = RGBColor(0x8B, 0x5C, 0xF6)   # NeuralFur
CARD = RGBColor(0x2A, 0x2A, 0x30)
GOOD = RGBColor(0x35, 0xC7, 0x59)
BAD  = RGBColor(0xE0, 0x5A, 0x5A)

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
ASSET = "/home/yyang/mnt/workspace"

def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def txt(s, x, y, w, h, lines, size=18, color=FG, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str): lines = [(lines, size, color, bold)]
    for i, ln in enumerate(lines):
        t, sz, c, b = ln if isinstance(ln, tuple) else (ln, size, color, bold)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = align
        run = p.add_run(); run.text = t; run.font.size = Pt(sz); run.font.bold = b
        run.font.color.rgb = c; run.font.name = "Arial"
    return tb

def box(s, x, y, w, h, text, fill=CARD, edge=None, tcolor=FG, size=15, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if edge: sp.line.color.rgb = edge; sp.line.width = Pt(1.5)
    else: sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for i, ln in enumerate(text.split("\n")):
        pp = p if i == 0 else tf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = ln; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = tcolor; r.font.name = "Arial"
    return sp

def arrow(s, x, y, w, color=SUB):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.32))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    return a

def img_fit(s, path, x, y, maxw, maxh, caption=None):
    im = Image.open(path); iw, ih = im.size; ar = iw/ih
    w = maxw; h = Emu(int(w / ar))
    if h > maxh: h = maxh; w = Emu(int(h * ar))
    px = x + (maxw - w)//2
    s.shapes.add_picture(path, px, y, w, h)
    if caption: txt(s, x, y + h + Pt(2), maxw, Inches(0.3), caption, size=11, color=SUB, align=PP_ALIGN.CENTER)
    return h

def header(s, title, accent):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, Inches(0.45), Inches(0.25), Inches(12.5), Inches(0.7), title, size=30, bold=True)

# ---------- prep a clean NeuralFur turntable frame ----------
os.makedirs(f"{ASSET}/_ppt_assets", exist_ok=True)
g = Image.open(f"{ASSET}/exps/neuralfur_final/00085-kotori_neuralfur_turntable.gif"); g.seek(8)
nf_frame = f"{ASSET}/_ppt_assets/nf_side.png"; g.convert("RGB").save(nf_frame)
# v9 render half (left), tight-crop to the dog (drop white margin)
v9 = Image.open(f"{ASSET}/exps/dog_lrm_fur_v9/it02500.png").convert("RGB")
v9l = v9.crop((0, 0, v9.width//2, v9.height))
import numpy as _np
_a = _np.asarray(v9l); _mask = (_a < 245).any(2); _ys, _xs = _np.where(_mask)
if len(_xs):
    _pad = 40
    v9l = v9l.crop((max(0, _xs.min()-_pad), max(0, _ys.min()-_pad),
                    min(v9l.width, _xs.max()+_pad), min(v9l.height, _ys.max()+_pad)))
v9_render = f"{ASSET}/_ppt_assets/v9_render.png"; v9l.save(v9_render)

# ================= SLIDE 1: TITLE =================
s = slide()
txt(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
    [("Dog-LRM Fur — Three Approaches", 40, FG, True)])
txt(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.4),
    [("Giving a 3DGS dog avatar fur that is decomposable (separable skin/undercoat vs. a", 18, SUB, False),
     ("simulatable fur layer) AND visually faithful.  Feed-forward · Two-stage · Cascade", 18, SUB, False)])
for i, (nm, cc) in enumerate([("A  Feed-forward", BLUE), ("B  Two-stage", RED), ("C  Cascade", GREEN)]):
    box(s, Inches(0.9 + i*4.0), Inches(5.5), Inches(3.6), Inches(0.7), nm, fill=cc, size=18)

# ================= helper: approach slide =================
def approach_slide(letter, name, accent, flow_fn, pros, cons, verdict, img, cap, imgw=Inches(4.3)):
    s = slide(); header(s, f"{letter})  {name}", accent)
    flow_fn(s)
    # pros/cons bottom-left
    y0 = Inches(5.55)
    txt(s, Inches(0.5), y0, Inches(8.0), Inches(1.6),
        [("✓  " + pros, 15, GOOD, True), ("✗  " + cons, 15, BAD, True),
         ("→  " + verdict, 15, FG, False)])
    # result bottom-right
    if img:
        img_fit(s, img, Inches(8.7), Inches(5.15), imgw, Inches(1.9), cap)
    return s

# ---- A flow ----
def flow_A(s):
    y = Inches(2.2); hh = Inches(0.95)
    box(s, Inches(0.5), y, Inches(1.8), hh, "Single\nimage", fill=CARD, edge=BLUE, size=15)
    arrow(s, Inches(2.4), y+Inches(0.32), Inches(0.7))
    box(s, Inches(3.2), y, Inches(2.2), hh, "DINOv2\nbackbone", fill=BLUE, size=15)
    arrow(s, Inches(5.5), y+Inches(0.32), Inches(0.7))
    box(s, Inches(6.35), Inches(1.55), Inches(3.0), Inches(0.85), "Body head → skin GS", fill=CARD, edge=BLUE, size=14)
    box(s, Inches(6.35), Inches(2.7), Inches(3.0), Inches(0.85), "Fur head → strands", fill=CARD, edge=BLUE, size=14)
    arrow(s, Inches(9.45), y+Inches(0.32), Inches(0.7))
    box(s, Inches(10.3), y, Inches(2.5), hh, "ONE forward\nbody + fur together", fill=BLUE, size=13)
    txt(s, Inches(6.35), Inches(3.75), Inches(6.3), Inches(0.5),
        "v9: + pixel-aligned splatter residual + adversarial  →  best texture", size=13, color=PURP, bold=True)

# ---- B flow ----
def flow_B(s):
    y = Inches(2.4); hh = Inches(0.95)
    box(s, Inches(0.5), y, Inches(2.8), hh, "Train body\n(furry photos)", fill=CARD, edge=RED, size=14)
    arrow(s, Inches(3.4), y+Inches(0.32), Inches(0.7), color=RED)
    box(s, Inches(4.2), y, Inches(2.3), hh, "❄  FREEZE", fill=RED, size=16)
    arrow(s, Inches(6.6), y+Inches(0.32), Inches(0.7), color=RED)
    box(s, Inches(7.45), y, Inches(2.6), hh, "Add fur layer\non top", fill=CARD, edge=RED, size=14)
    arrow(s, Inches(10.15), y+Inches(0.32), Inches(0.7), color=RED)
    box(s, Inches(11.0), y, Inches(1.9), hh, "✗ redundant", fill=RED, size=13)
    txt(s, Inches(0.5), Inches(3.7), Inches(12), Inches(0.9),
        [("Body was trained on furry photos → coat already BAKED into body.", 15, SUB, False),
         ("Frozen body already looks right → the fur has nowhere to live → adding it RAISES the error.", 15, SUB, False)])

# ---- C flow ----
def flow_C(s):
    y = Inches(1.9); hh = Inches(0.95)
    box(s, Inches(1.0), y, Inches(3.4), hh, "Body gaussians\n(recede → dark undercoat)", fill=CARD, edge=GREEN, size=14)
    box(s, Inches(1.0), Inches(3.45), Inches(3.4), hh, "Fur strands\n(semi-transparent, carry coat)", fill=CARD, edge=GREEN, size=14)
    # coupling double arrow (native vertical, in the gap between the two boxes)
    da = s.shapes.add_shape(MSO_SHAPE.UP_DOWN_ARROW, Inches(2.45), Inches(2.9), Inches(0.5), Inches(0.5))
    da.fill.solid(); da.fill.fore_color.rgb = GREEN; da.line.fill.background(); da.shadow.inherit = False
    txt(s, Inches(4.5), Inches(2.55), Inches(2.2), Inches(0.6), "coupled\n(recession)", size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    arrow(s, Inches(6.7), Inches(2.85), Inches(0.7), color=GREEN)
    box(s, Inches(7.55), Inches(2.15), Inches(2.5), Inches(0.85), "composite ≈ GT", fill=GREEN, size=14)
    box(s, Inches(7.55), Inches(3.35), Inches(2.5), Inches(0.85), "sway reveals\nundercoat", fill=CARD, edge=GREEN, size=13)
    box(s, Inches(10.3), Inches(2.75), Inches(2.5), Inches(0.9), "DECOMPOSABLE\n+ SIMULATABLE", fill=GREEN, size=13)
    txt(s, Inches(1.0), Inches(4.55), Inches(9), Inches(0.4),
        "per-scene multi-view optimization (ceiling probe) · face excluded via w_face", size=12, color=SUB)

approach_slide("A", "Feed-forward (joint)   v6 → v9", BLUE, flow_A,
    "Best texture; the ONLY one that generalizes (single image → fur)",
    "Coat baked into body → NOT decomposable; face needs a nofur mask",
    "Great looking fur, but no movable/separable fur layer",
    v9_render, "v9 feed-forward (render): sharpest texture")

approach_slide("B", "Two-stage (frozen)", RED, flow_B,
    "Clean modularity in principle",
    "Frozen body already baked the coat → fur redundant → adding fur HURTS L1",
    "Dead end — decoupling without letting skin recede makes fur useless",
    None, None)
# B: add a 'no usable result' card
sB = prs.slides[-1]
box(sB, Inches(9.0), Inches(5.2), Inches(3.6), Inches(1.5), "no usable result\n(dead end)", fill=CARD, edge=RED, tcolor=SUB, size=15)

approach_slide("C", "Cascade (coupled + recession)   v11", GREEN, flow_C,
    "Decomposable + simulatable; composite ≈ GT",
    "Per-scene optimization (not feed-forward); texture softer than v9",
    "The route that actually separates skin ↔ fur",
    f"{ASSET}/exps/fur_v11_kotori_clean/00085-kotori_decomp.png",
    "cascade decomposition: skin → undercoat → fur → composite → GT", imgw=Inches(4.4))

# ================= COMPARISON TABLE =================
s = slide(); header(s, "Axis-by-axis comparison", RGBColor(0x9C,0x9C,0xA6))
rows = [
    ("Axis", "A) Feed-forward", "B) Two-stage", "C) Cascade"),
    ("Skin ↔ fur", "parallel heads, no coupling", "frozen, hard-decoupled", "coupled + recession"),
    ("Appearance lives in", "body (coat baked)", "body (coat baked)", "split: undercoat + fur"),
    ("Decomposable?", "✗", "✗", "✓"),
    ("Role of fur", "decorative", "redundant", "load-bearing (carries coat)"),
    ("Simulatable", "weak", "weak", "✓ (sway works)"),
    ("Texture", "best (v9)", "n/a", "softer; needs sharpening"),
    ("1-image generalization", "✓", "partial", "✗ (per-scene)"),
    ("Static composite L1", "~body-shell (low)", "fur raises it", "0.010–0.015"),
    ("Face", "nofur mask", "—", "w_face excluded"),
]
tb = s.shapes.add_table(len(rows), 4, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.1)).table
tb.columns[0].width = Inches(2.7)
for c in range(1, 4): tb.columns[c].width = Inches(3.2)
accents = [None, BLUE, RED, GREEN]
for ri, row in enumerate(rows):
    tb.rows[ri].height = Inches(0.47)
    for ci, val in enumerate(row):
        cell = tb.cell(ri, ci); cell.margin_left = Pt(6); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        cell.fill.solid()
        if ri == 0: cell.fill.fore_color.rgb = accents[ci] or RGBColor(0x33,0x33,0x3A)
        else: cell.fill.fore_color.rgb = CARD if ri % 2 else RGBColor(0x24,0x24,0x2A)
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
        r.font.size = Pt(12.5 if ri else 13); r.font.name = "Arial"
        r.font.bold = (ri == 0 or ci == 0 or (ci == 3 and val in ("✓","coupled + recession","split: undercoat + fur","load-bearing (carries coat)","✓ (sway works)")))
        r.font.color.rgb = FG if (ri == 0) else (SUB if ci == 0 else FG)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
txt(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.8),
    [("Disproven myth: structured/simulatable fur cannot beat a smooth body-shell on static L1 →  L1<0.01 is a body-shell metric.", 13, RGBColor(0xE7,0xC0,0x6A), True),
     ("Fur's value = decomposability + dynamics + perceptual sharpness, NOT static L1.", 13, SUB, False)])

# ================= RESULTS COMPARISON =================
s = slide(); header(s, "Best results", RGBColor(0x9C,0x9C,0xA6))
cols = [
    ("A  Feed-forward (v9)", BLUE, v9_render, "sharpest texture · but coat baked in"),
    ("C  Cascade (kotori)", GREEN, f"{ASSET}/exps/fur_v11_kotori_clean/cmp.png", "4 views · decomposable"),
    ("NeuralFur geom + our colour", PURP, f"{ASSET}/exps/neuralfur_final/00085-kotori_neuralfur_colored.png", "GH strand geometry, recoloured (no speckle)"),
]
cw = Inches(4.05)
for i, (nm, cc, img, cap) in enumerate(cols):
    x = Inches(0.5 + i*4.28)
    box(s, x, Inches(1.2), cw, Inches(0.55), nm, fill=cc, size=15)
    img_fit(s, img, x, Inches(1.95), cw, Inches(4.2), cap)
txt(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
    "B) Two-stage produced no usable result (dead end).", size=13, color=SUB)

# ================= RECOMMENDATION =================
s = slide(); header(s, "Where this points", PURP)
txt(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "Best texture = A (v9).   Best decomposition = C (cascade).   They have not yet been combined.", size=18, bold=True)
items = [
    ("1  Combine v9 texture + cascade decomposition", "Use v9's fur head as the Stage-2 layer, rooted on Stage-1 bare skin; add recession so the body becomes true undercoat; exclude the face. → v9 texture + decomposable + bald face."),
    ("2  Feed-forward the cascade (deploy)", "Train an image-conditioned head to predict {fur op/colour/geometry + body recession} from ONE image, supervised by the per-scene cascade results. → A's generalization + C's decomposability."),
    ("3  NeuralFur as a geometry source", "This round proved 'external strand geometry + our colour-query' works cleanly — a drop-in quality source for the cascade's fur geometry."),
]
y = Inches(2.1)
for t, d in items:
    box(s, Inches(0.6), y, Inches(0.5), Inches(0.5), "", fill=PURP, size=1)
    txt(s, Inches(1.25), y-Inches(0.05), Inches(11.4), Inches(1.2),
        [(t, 17, FG, True), (d, 14, SUB, False)])
    y += Inches(1.55)

out = f"{ASSET}/Dog_Fur_Three_Approaches.pptx"
prs.save(out); print("saved", out, "slides", len(prs.slides._sldIdLst))
